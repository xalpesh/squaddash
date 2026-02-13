from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# --- MCKINSEY COLORS ---
MCK_NAVY = RGBColor(16, 55, 92)
MCK_BLUE = RGBColor(0, 112, 192)
MCK_CARD_BG = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(64, 64, 64)

def create_sow_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank Layout

    # 1. ACTION TITLE & HEADER (SOW Summary)
    # The SOW emphasizes establishing workload trends and rolling forecasts
    draw_header(slide, 
                "Pillar 1: Demand Forecasting", 
                "Predictive 30/60/90-day forecasts drive skill readiness vs. reactive hiring.")

    # 2. LEFT COLUMN: METRICS CARD (SOW Specifics)
    # SOW Ref: "Collect and maintain historical data... predict 30/60/90/180 day rolling forecast"
    draw_metric_card(slide, 
                     focus="Trend Analysis & Roadmap Review", 
                     metrics=["• 30/60/90/180 Day Rolling Forecast", 
                              "• Ticket Trend Volume (BAU)", 
                              "• Skill-wise Gap Analysis"],
                     role="Supplier Portfolio Lead")

    # 3. RIGHT COLUMN: PHASE MAP VISUAL
    # "Map on what phases of project this data will be collected"
    # We map the SOW inputs (Historical Data, Roadmap, Forecast) to phases.
    draw_phase_map(slide)

    # 4. KICKER
    draw_kicker(slide, "Data-driven forecasting ensures the right skills are available before the project kick-off.")

    prs.save('Pillar1_SOW_Aligned.pptx')
    print("Slide Generated: Pillar1_SOW_Aligned.pptx")

# --- DRAWING FUNCTIONS ---

def draw_header(slide, title, sub):
    # Main Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = MCK_NAVY
    
    # Blue Separator
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = MCK_BLUE
    line.line.fill.background()
    
    # Subtitle (Summary from SOW)
    tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = sub
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = MCK_NAVY

def draw_metric_card(slide, focus, metrics, role):
    # Dark Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(3.5), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = MCK_NAVY
    header.line.color.rgb = MCK_NAVY
    header.text_frame.text = "Focus Area / SOW Metrics"
    header.text_frame.paragraphs[0].font.color.rgb = WHITE
    header.text_frame.paragraphs[0].font.bold = True
    
    # Light Body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.8), Inches(3.5), Inches(3.5))
    body.fill.solid()
    body.fill.fore_color.rgb = MCK_CARD_BG
    body.line.color.rgb = MCK_CARD_BG
    
    tf = body.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    
    # Content
    add_text(tf, "KEY ACTIVITY:", focus, True)
    
    # Metrics Loop
    p_label = tf.add_paragraph()
    p_label.text = "MEASURABLE METRICS:"
    p_label.font.size = Pt(11)
    p_label.font.bold = True
    p_label.font.color.rgb = MCK_NAVY
    
    for m in metrics:
        p_val = tf.add_paragraph()
        p_val.text = m
        p_val.font.size = Pt(12)
        p_val.font.color.rgb = DARK_TEXT
        p_val.space_after = Pt(2)
    
    p_spacer = tf.add_paragraph()
    p_spacer.space_after = Pt(10)

    add_text(tf, "PRIMARY OWNER:", role, True)

def add_text(tf, label, value, bold_label=False):
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = MCK_NAVY
    
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_TEXT
    p2.space_after = Pt(14)

def draw_phase_map(slide):
    # Header for the Visual
    lbl = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(6), Inches(0.5))
    lbl.text_frame.text = "Data Collection Phase Map (Aligned to SDLC)"
    lbl.text_frame.paragraphs[0].font.bold = True
    lbl.text_frame.paragraphs[0].font.color.rgb = MCK_NAVY
    lbl.text_frame.paragraphs[0].font.size = Pt(14)

    # 3 Phases: Continuous -> Planning -> Initiation
    phases = [
        {"name": "Continuous / BAU", "data": "Historical Ticket Trends\n& Workload Data"},
        {"name": "Strategic Planning", "data": "Roadmap Review\n& Project Pipeline"},
        {"name": "Project Initiation", "data": "30/60/90 Day\nRolling Forecast"}
    ]
    
    left_x = Inches(4.5)
    top_y = Inches(2.5)
    
    for i, phase in enumerate(phases):
        # Draw Chevron or Arrow Box
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left_x, top_y, Inches(2.6), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MCK_BLUE
        shape.line.color.rgb = WHITE
        shape.line.width = Pt(2)
        
        # Phase Title (Top of Chevron)
        tf = shape.text_frame
        tf.text = phase["name"] + "\n\n\n" + phase["data"]
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow Connector (if not last)
        if i < 2:
            conn = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_x + Inches(2.4), top_y + Inches(1.4), Inches(0.5), Inches(0.2))
            conn.fill.solid()
            conn.fill.fore_color.rgb = MCK_NAVY
            
        left_x += Inches(2.9)
    
    # Legend/Note
    note = slide.shapes.add_textbox(Inches(4.5), Inches(5.8), Inches(8), Inches(0.5))
    note.text_frame.text = "* Historical data establishes baseline; Roadmap review predicts peaks."
    note.text_frame.paragraphs[0].font.size = Pt(10)
    note.text_frame.paragraphs[0].font.italic = True

def draw_kicker(slide, text):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = MCK_CARD_BG
    box.line.fill.background()
    
    kp = box.text_frame.paragraphs[0]
    kp.text = text
    kp.font.size = Pt(14)
    kp.font.italic = True
    kp.font.color.rgb = MCK_NAVY
    kp.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_sow_slide()